"""Document parsers - extract text from various file formats."""
from __future__ import annotations

import csv
import io
import json
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Any


@dataclass
class ParsedContent:
    """Result of parsing a document."""
    text: str
    meta: dict[str, Any] = field(default_factory=dict)
    html: str | None = None
    sections: list[dict[str, Any]] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)


class Parser(ABC):
    """Base class for document parsers."""

    @abstractmethod
    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        """Parse document content and extract text."""
        pass

    @abstractmethod
    def can_parse(self, mime_type: str, filename: str) -> bool:
        """Check if this parser can handle the given file."""
        pass


def _get_extension(filename: str) -> str:
    """Get lowercase file extension with dot."""
    if "." in filename:
        return "." + filename.rsplit(".", 1)[-1].lower()
    return ""


# ─── Text / Markdown ─────────────────────────────────────────────────────────


class TextParser(Parser):
    """Parser for plain text files."""

    SUPPORTED_MIMES = {"text/plain", "text/x-log"}
    SUPPORTED_EXTENSIONS = {".txt", ".log", ".rst", ".ini", ".cfg", ".conf", ".env"}

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type in self.SUPPORTED_MIMES:
            return True
        return _get_extension(filename) in self.SUPPORTED_EXTENSIONS

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")
        return ParsedContent(text=text.strip(), meta={"encoding": "utf-8"})


class MarkdownParser(Parser):
    """Parser for Markdown files with heading structure extraction."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "text/markdown":
            return True
        return _get_extension(filename) in {".md", ".markdown"}

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        sections = []
        current_heading_path: list[str] = []

        for line in text.split("\n"):
            match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if match:
                level = len(match.group(1))
                heading = match.group(2).strip()
                while len(current_heading_path) >= level:
                    current_heading_path.pop()
                current_heading_path.append(heading)
                sections.append({
                    "heading": heading,
                    "level": level,
                    "path": list(current_heading_path),
                })

        return ParsedContent(text=text.strip(), meta={"format": "markdown"}, sections=sections)


# ─── PDF ─────────────────────────────────────────────────────────────────────


class PDFParser(Parser):
    """Parser for PDF files using PyPDF2 or pdfminer."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "application/pdf":
            return True
        return _get_extension(filename) == ".pdf"

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        text_parts: list[str] = []
        meta: dict[str, Any] = {"format": "pdf", "pages": 0}

        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader  # type: ignore[no-redef]
            except ImportError:
                try:
                    from pdfminer.high_level import extract_text
                    text = extract_text(io.BytesIO(content))
                    return ParsedContent(text=text.strip(), meta={"format": "pdf", "parser": "pdfminer"})
                except ImportError:
                    raise RuntimeError("No PDF parser available. Install pypdf or pdfminer.six")

        reader = PdfReader(io.BytesIO(content))
        meta["pages"] = len(reader.pages)

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

        return ParsedContent(text="\n\n".join(text_parts).strip(), meta=meta)


# ─── DOCX ────────────────────────────────────────────────────────────────────


class DocxParser(Parser):
    """Parser for Word documents using python-docx."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return True
        return _get_extension(filename) == ".docx"

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError("python-docx not installed. Run: pip install python-docx")

        doc = Document(io.BytesIO(content))
        text_parts: list[str] = []
        sections: list[dict[str, Any]] = []
        current_heading_path: list[str] = []

        for para in doc.paragraphs:
            style_name = para.style.name if para.style else ""

            # Detect headings
            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.replace("Heading", "").strip())
                except ValueError:
                    level = 1
                heading = para.text.strip()
                if heading:
                    while len(current_heading_path) >= level:
                        current_heading_path.pop()
                    current_heading_path.append(heading)
                    sections.append({
                        "heading": heading,
                        "level": level,
                        "path": list(current_heading_path),
                    })

            if para.text.strip():
                text_parts.append(para.text)

        # Also extract tables
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            table_data: list[list[str]] = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables.append(table_data)
                # Convert table to markdown text
                md_table = _table_to_markdown(table_data)
                text_parts.append(f"\n[Table]\n{md_table}\n")

        return ParsedContent(
            text="\n".join(text_parts).strip(),
            meta={"format": "docx", "paragraphs": len(doc.paragraphs), "tables": len(tables)},
            sections=sections,
            tables=tables,
        )


# ─── XLSX ────────────────────────────────────────────────────────────────────


class XlsxParser(Parser):
    """Parser for Excel files using openpyxl."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return True
        return _get_extension(filename) in {".xlsx", ".xls"}

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"## Sheet: {sheet_name}\n")

            table_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):
                    table_data.append(row_data)

            if table_data:
                tables.append(table_data)
                md_table = _table_to_markdown(table_data)
                text_parts.append(md_table)

        wb.close()

        return ParsedContent(
            text="\n".join(text_parts).strip(),
            meta={"format": "xlsx", "sheets": wb.sheetnames, "tables": len(tables)},
            tables=tables,
        )


# ─── PPTX ────────────────────────────────────────────────────────────────────


class PptxParser(Parser):
    """Parser for PowerPoint files using python-pptx."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return True
        return _get_extension(filename) == ".pptx"

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(io.BytesIO(content))
        text_parts: list[str] = []
        sections: list[dict[str, Any]] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            title_text = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                            # First text is likely the title
                            if not title_text and shape.shape_type is not None:
                                title_text = text

                if shape.has_table:
                    table_data: list[list[str]] = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        md_table = _table_to_markdown(table_data)
                        slide_texts.append(f"\n[Table]\n{md_table}\n")

            if slide_texts:
                heading = title_text or f"Slide {i}"
                sections.append({
                    "heading": heading,
                    "level": 2,
                    "path": [heading],
                })
                text_parts.append(f"## Slide {i}: {title_text or '(no title)'}\n")
                text_parts.extend(slide_texts)

        return ParsedContent(
            text="\n".join(text_parts).strip(),
            meta={"format": "pptx", "slides": len(prs.slides)},
            sections=sections,
        )


# ─── HTML ────────────────────────────────────────────────────────────────────


class HTMLParser(Parser):
    """Parser for HTML files. Strips tags, preserves text."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type in {"text/html", "application/xhtml+xml"}:
            return True
        return _get_extension(filename) in {".html", ".htm"}

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            html_text = content.decode("utf-8")
        except UnicodeDecodeError:
            html_text = content.decode("latin-1")

        # Try readability for article extraction
        try:
            from readability import Document as ReadabilityDocument
            doc = ReadabilityDocument(html_text)
            text = doc.summary()
            # Strip HTML tags from summary
            text = _strip_html_tags(text)
            title = doc.title()
            return ParsedContent(
                text=text.strip(),
                html=html_text,
                meta={"format": "html", "title": title, "parser": "readability"},
            )
        except ImportError:
            pass

        # Fallback: simple tag stripping
        text = _strip_html_tags(html_text)
        # Extract title
        title_match = re.search(r"<title[^>]*>([^<]+)</title>", html_text, re.IGNORECASE)
        title = title_match.group(1).strip() if title_match else ""

        return ParsedContent(
            text=text.strip(),
            html=html_text,
            meta={"format": "html", "title": title},
        )


# ─── CSV ─────────────────────────────────────────────────────────────────────


class CSVParser(Parser):
    """Parser for CSV files. Converts to markdown table."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "text/csv":
            return True
        return _get_extension(filename) == ".csv"

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        reader = csv.reader(io.StringIO(text))
        rows = list(reader)

        if not rows:
            return ParsedContent(text="", meta={"format": "csv", "rows": 0})

        md_table = _table_to_markdown(rows)

        return ParsedContent(
            text=md_table,
            meta={"format": "csv", "rows": len(rows), "columns": len(rows[0]) if rows else 0},
            tables=[rows],
        )


# ─── Jupyter Notebook ────────────────────────────────────────────────────────


class NotebookParser(Parser):
    """Parser for Jupyter notebooks (.ipynb)."""

    def can_parse(self, mime_type: str, filename: str) -> bool:
        if mime_type == "application/x-ipynb+json":
            return True
        return _get_extension(filename) == ".ipynb"

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            nb = json.loads(content.decode("utf-8"))
        except (json.JSONDecodeError, UnicodeDecodeError) as e:
            raise ValueError(f"Invalid notebook JSON: {e}")

        text_parts: list[str] = []
        sections: list[dict[str, Any]] = []
        current_heading_path: list[str] = []

        cells = nb.get("cells", [])
        for i, cell in enumerate(cells):
            cell_type = cell.get("cell_type", "")
            source = cell.get("source", [])
            if isinstance(source, list):
                source_text = "".join(source)
            else:
                source_text = str(source)

            if cell_type == "markdown":
                text_parts.append(source_text)
                # Extract headings
                for line in source_text.split("\n"):
                    match = re.match(r"^(#{1,6})\s+(.+)$", line)
                    if match:
                        level = len(match.group(1))
                        heading = match.group(2).strip()
                        while len(current_heading_path) >= level:
                            current_heading_path.pop()
                        current_heading_path.append(heading)
                        sections.append({
                            "heading": heading,
                            "level": level,
                            "path": list(current_heading_path),
                        })

            elif cell_type == "code":
                # Include code with language tag
                text_parts.append(f"```python\n{source_text}\n```")

                # Include outputs
                outputs = cell.get("outputs", [])
                for output in outputs:
                    output_type = output.get("output_type", "")
                    if output_type == "stream":
                        stream_text = output.get("text", [])
                        if isinstance(stream_text, list):
                            stream_text = "".join(stream_text)
                        if stream_text.strip():
                            text_parts.append(f"Output:\n{stream_text}")
                    elif output_type in ("execute_result", "display_data"):
                        data = output.get("data", {})
                        if "text/plain" in data:
                            plain = data["text/plain"]
                            if isinstance(plain, list):
                                plain = "".join(plain)
                            text_parts.append(f"Result:\n{plain}")

        return ParsedContent(
            text="\n\n".join(text_parts).strip(),
            meta={"format": "jupyter", "cells": len(cells)},
            sections=sections,
        )


# ─── Code Files ──────────────────────────────────────────────────────────────


class CodeParser(Parser):
    """Parser for source code files. Preserves structure."""

    LANGUAGE_MAP = {
        ".py": "python", ".pyi": "python", ".pyx": "python",
        ".js": "javascript", ".mjs": "javascript", ".cjs": "javascript",
        ".ts": "typescript", ".tsx": "typescript",
        ".jsx": "javascript",
        ".java": "java",
        ".kt": "kotlin", ".kts": "kotlin",
        ".go": "go",
        ".rs": "rust",
        ".rb": "ruby",
        ".php": "php",
        ".swift": "swift",
        ".c": "c", ".h": "c",
        ".cpp": "cpp", ".cxx": "cpp", ".cc": "cpp", ".hpp": "cpp", ".hxx": "cpp",
        ".cs": "csharp",
        ".scala": "scala",
        ".r": "r", ".R": "r",
        ".sh": "shell", ".bash": "shell", ".zsh": "shell",
        ".sql": "sql",
        ".yaml": "yaml", ".yml": "yaml",
        ".json": "json",
        ".xml": "xml",
        ".toml": "toml",
        ".ini": "ini", ".cfg": "ini",
        ".dart": "dart",
        ".lua": "lua",
        ".pl": "perl", ".pm": "perl",
        ".hs": "haskell",
        ".clj": "clojure", ".cljs": "clojure",
        ".ex": "elixir", ".exs": "elixir",
        ".erl": "erlang",
        ".vue": "vue",
        ".svelte": "svelte",
    }

    SUPPORTED_EXTENSIONS = set(LANGUAGE_MAP.keys())

    def can_parse(self, mime_type: str, filename: str) -> bool:
        ext = _get_extension(filename)
        return ext in self.SUPPORTED_EXTENSIONS

    def parse(self, content: bytes, mime_type: str, filename: str) -> ParsedContent:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1")

        ext = _get_extension(filename)
        language = self.LANGUAGE_MAP.get(ext, "text")

        # Extract function/class definitions as sections
        sections = _extract_code_structure(text, language)

        # Wrap in code block
        code_text = f"```{language}\n{text}\n```"

        return ParsedContent(
            text=code_text,
            meta={"format": "code", "language": language, "lines": text.count("\n") + 1},
            sections=sections,
        )


# ─── Helpers ─────────────────────────────────────────────────────────────────


def _strip_html_tags(html: str) -> str:
    """Strip HTML tags and decode entities."""
    import html as html_module
    # Remove script and style blocks
    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.IGNORECASE)
    # Replace block elements with newlines
    text = re.sub(r"<(br|p|div|h[1-6]|li|tr)[^>]*>", "\n", text, flags=re.IGNORECASE)
    # Strip remaining tags
    text = re.sub(r"<[^>]+>", "", text)
    # Decode HTML entities
    text = html_module.unescape(text)
    # Collapse whitespace
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _table_to_markdown(rows: list[list[str]]) -> str:
    """Convert a table (list of rows) to markdown format."""
    if not rows:
        return ""

    # Normalize column count
    max_cols = max(len(row) for row in rows)
    normalized = []
    for row in rows:
        padded = row + [""] * (max_cols - len(row))
        normalized.append(padded)

    # Build markdown
    lines: list[str] = []
    # Header
    lines.append("| " + " | ".join(normalized[0]) + " |")
    lines.append("| " + " | ".join("---" for _ in normalized[0]) + " |")
    # Body
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _extract_code_structure(text: str, language: str) -> list[dict[str, Any]]:
    """Extract function/class definitions from code as sections."""
    sections: list[dict[str, Any]] = []

    patterns: list[tuple[str, str]] = []
    if language == "python":
        patterns = [
            (r"^class\s+(\w+)", "class"),
            (r"^\s*def\s+(\w+)", "function"),
            (r"^\s*async\s+def\s+(\w+)", "function"),
        ]
    elif language in ("javascript", "typescript", "java", "csharp", "go", "rust"):
        patterns = [
            (r"(?:class|interface|struct|enum)\s+(\w+)", "class"),
            (r"(?:function|def|fn|func)\s+(\w+)", "function"),
            (r"(?:const|let|var)\s+(\w+)\s*=\s*(?:async\s+)?(?:\([^)]*\)|[^=])\s*=>", "function"),
        ]

    for line_no, line in enumerate(text.split("\n"), 1):
        for pattern, kind in patterns:
            match = re.match(pattern, line)
            if match:
                name = match.group(1)
                sections.append({
                    "heading": f"{kind}: {name}",
                    "level": 3 if kind == "class" else 4,
                    "line": line_no,
                    "kind": kind,
                    "name": name,
                })
                break

    return sections


# ─── Registry ────────────────────────────────────────────────────────────────

# Ordered by specificity (most specific first)
_PARSERS: list[Parser] = [
    PDFParser(),
    DocxParser(),
    XlsxParser(),
    PptxParser(),
    NotebookParser(),
    HTMLParser(),
    CSVParser(),
    MarkdownParser(),
    CodeParser(),
    TextParser(),  # Fallback
]


def get_parser(mime_type: str, filename: str) -> Parser | None:
    """Get the appropriate parser for a file."""
    for parser in _PARSERS:
        if parser.can_parse(mime_type, filename):
            return parser
    return None


def parse_document(content: bytes, mime_type: str, filename: str) -> ParsedContent:
    """Parse a document and extract text."""
    parser = get_parser(mime_type, filename)
    if parser is None:
        raise ValueError(f"No parser available for {mime_type} ({filename})")
    return parser.parse(content, mime_type, filename)


def supported_extensions() -> set[str]:
    """Return all supported file extensions."""
    exts: set[str] = set()
    for parser in _PARSERS:
        if hasattr(parser, "SUPPORTED_EXTENSIONS"):
            exts.update(parser.SUPPORTED_EXTENSIONS)  # type: ignore[attr-defined]
        if hasattr(parser, "LANGUAGE_MAP"):
            exts.update(parser.LANGUAGE_MAP.keys())  # type: ignore[attr-defined]
    return exts
